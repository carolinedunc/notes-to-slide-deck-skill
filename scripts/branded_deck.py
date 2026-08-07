from pptx import Presentation 
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1B, 0x3A, 0x5C)
ORANGE = RGBColor(0xE8, 0x79, 0x2C)
GOLD = RGBColor(0xF2, 0xB2, 0x33)
BACKGROUND = RGBColor(0xFA, 0xFA, 0xFA)
BLACK = RGBColor(0x2B, 0x2B, 0x2B)
GREY = RGBColor(0x80, 0x80, 0x80)

def add_brand_slide(prs, title_text, slide_num):
      slide = prs.slides.add_slide(prs.slide_layouts[6])

      background = slide.background
      fill = background.fill
      fill.solid()
      fill.fore_color.rgb = BACKGROUND
      
      #wordmark formatting 
      wordmark_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.15), Inches(3), Inches(0.3))
      wordmarkFrame = wordmark_textbox.text_frame
      wordmarkFrame.text = "DUNCAN & CO."
      wordmarkRun = wordmarkFrame.paragraphs[0].runs[0]
      wordmarkRun.font.size = Pt(12)
      wordmarkRun.font.bold = True
      wordmarkRun.font.color.rgb = NAVY

      #title formatting
      title_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(8.5), Inches(1.0))
      textframe = title_textbox.text_frame 
      textframe.word_wrap = True
      textframe.text = title_text

      run = textframe.paragraphs[0].runs[0]
      if len(title_text) > 45:
            run.font.size = Pt(20)
      else:
           run.font.size = Pt(25)
      run.font.bold = True 
      run.font.color.rgb = NAVY

      #orange line after title 
      line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(8.5), Inches(0.02))
      line1.fill.solid()
      line1.fill.fore_color.rgb = ORANGE
      line1.line.fill.background()

      #footer formatting
      footer1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.3), Inches(8.5), Inches(0.02))
      footer1.fill.solid()
      footer1.fill.fore_color.rgb = GREY
      footer1.line.fill.background()

      #footer wordmark
      footer_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(4), Inches(0.3))
      footerFrame = footer_textbox.text_frame
      footerFrame.text = "Duncan & Co."
      footerRun = footerFrame.paragraphs[0].runs[0]
      footerRun.font.size = Pt(9)
      footerRun.font.color.rgb = GREY

      #slide number
      number_textbox = slide.shapes.add_textbox(Inches(8.5), Inches(5.1), Inches(0.5), Inches(0.2))
      numberFrame = number_textbox.text_frame
      numberFrame.text = str(slide_num)
      numberRun = numberFrame.paragraphs[0].runs[0]
      numberRun.font.size = Pt(9)
      numberRun.font.color.rgb = GREY
      numberFrame.paragraphs[0].alignment = PP_ALIGN.RIGHT

      #logo placement
      logo = slide.shapes.add_picture("../assets/logo.png", Inches(0.25), Inches(0.15), height=Inches(0.7))

      return slide

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.5)

    add_brand_slide(prs, "Q1'26 Report", 1)
    add_brand_slide(prs, "Key Themes & Overall Takeaways", 2)
    add_brand_slide(prs, "Risk", 3)  # very short title
    add_brand_slide(prs, "Customer Pipeline and Revenue Guidance for the Upcoming Fiscal Year", 4)  # very long title
    add_brand_slide(prs, "Open Questions & Follow-Ups from Site Visit", 5)  # medium-long

    prs.save("../example/template_test.pptx")